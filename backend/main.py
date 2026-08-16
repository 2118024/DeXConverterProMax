import logging
import shutil
import uuid
import asyncio
import os
import subprocess
import zipfile
from pathlib import Path
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware

import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageEnhance

SOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"

app = FastAPI(title="Universal Document Converter Pro", version="5.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

TEMP_DIR = Path("temp_files")
TEMP_DIR.mkdir(exist_ok=True)

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(name)s: %(message)s")
logger = logging.getLogger("converter")

def cleanup_file(path: Path):
    """Menghapus file atau direktori sementara setelah direspons ke client."""
    try:
        if path.exists():
            if path.is_file():
                path.unlink()
            elif path.is_dir():
                shutil.rmtree(path)
            logger.info(f"Berhasil membersihkan: {path}")
    except Exception as e:
        logger.error(f"Gagal membersihkan {path}: {e}")

def convert_image_to_hd_pdf(input_path: Path, output_path: Path):
    """Mengonversi foto/gambar ke PDF dengan peningkatan kualitas HD, ketajaman, & kontras jernih."""
    image = Image.open(input_path)
    
    if image.mode != "RGB":
        image = image.convert("RGB")
        
    enhancer = ImageEnhance.Sharpness(image)
    image_sharp = enhancer.enhance(2.0)
    
    contrast_enhancer = ImageEnhance.Contrast(image_sharp)
    image_final = contrast_enhancer.enhance(1.3)
    
    # Simpan sebagai PDF resolusi tinggi (300 DPI)
    image_final.save(output_path, "PDF", resolution=300.0, save_all=True)
    return output_path

def run_libreoffice_conversion(input_path: Path, output_dir: Path, target_ext: str):
    """Fungsi standar menggunakan LibreOffice CLI dengan isolasi Environment agar tidak bentrok."""
    env = os.environ.copy()
    env.pop("PYTHONHOME", None)
    env.pop("PYTHONPATH", None)
    
    cmd = [
        SOFFICE_PATH, 
        "--headless", 
        "--convert-to", target_ext, 
        "--outdir", str(output_dir), 
        str(input_path)
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True, env=env)
    
    if result.returncode != 0:
        err_msg = result.stderr.strip() or result.stdout.strip()
        raise RuntimeError(f"LibreOffice Error: {err_msg}")
        
    output_file = output_dir / f"{input_path.stem}.{target_ext}"
    if not output_file.exists():
        raise FileNotFoundError("File hasil konversi tidak ditemukan di direktori output.")
        
    return output_file

def run_smart_conversion_with_template(input_path: Path, output_dir: Path, target_ext: str, template_style: str):
    ext_clean = target_ext.replace(".", "").lower()
    input_ext = input_path.suffix.lower()
    output_file = output_dir / f"{input_path.stem}.{ext_clean}"
    if input_ext in [".jpg", ".jpeg", ".png", ".webp"] and ext_clean == "pdf":
        return convert_image_to_hd_pdf(input_path, output_file)
    if input_ext in [".docx", ".doc"] and ext_clean == "pptx":
        doc = Document(input_path)
        prs = Presentation()
        
        if template_style == "corporate":
            title_color = RGBColor(15, 23, 42)    # Dark Navy
        elif template_style == "creative":
            title_color = RGBColor(225, 29, 72)   # Rose Vibrant
        else: # modern
            title_color = RGBColor(37, 99, 235)   # Primary Blue
        
        slide_layout = prs.slide_layouts[0] 
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = input_path.stem.replace("_", " ").title()
        if title.text_frame.paragraphs:
            title.text_frame.paragraphs[0].font.color.rgb = title_color
            
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Dikonversi otomatis dengan tema: {template_style.capitalize()}"
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text and len(text) > 3:
                bullet_layout = prs.slide_layouts[1] 
                slide_content = prs.slides.add_slide(bullet_layout)
                slide_content.shapes.title.text = "Poin Dokumen"
                if len(slide_content.placeholders) > 1:
                    slide_content.placeholders[1].text = text
                    
        prs.save(output_file)
        return output_file
    elif input_ext in [".docx", ".doc"] and ext_clean == "xlsx":
        doc = Document(input_path)
        data = []
        if doc.tables:
            for table in doc.tables:
                for row in table.rows:
                    data.append([cell.text.strip() for cell in row.cells])
        else:
            for para in doc.paragraphs:
                if para.text.strip():
                    data.append([para.text.strip()])
                    
        df = pd.DataFrame(data)
        df.to_excel(output_file, index=False, header=False)
        return output_file
    elif input_ext in [".xlsx", ".xls"] and ext_clean == "docx":
        df = pd.read_excel(input_path)
        doc = Document()
        doc.add_heading(f"Laporan Data: {input_path.stem}", level=1)
        
        table = doc.add_table(rows=df.shape[0] + 1, cols=df.shape[1])
        for j, col_name in enumerate(df.columns):
            table.cell(0, j).text = str(col_name)
        for i, row in df.iterrows():
            for j, val in enumerate(row):
                table.cell(i + 1, j).text = str(val) if pd.notna(val) else ""
                
        doc.save(output_file)
        return output_file
    else:
        return run_libreoffice_conversion(input_path, output_dir, ext_clean)

@app.post("/api/convert")
async def convert_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    target_format: str = Form(...),
    template_style: str = Form("modern")
):
    file_id = uuid.uuid4().hex
    input_path = TEMP_DIR / f"{file_id}_{file.filename}"
    
    try:
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal menyimpan file: {e}")

    try:
        output_path = await asyncio.to_thread(
            run_smart_conversion_with_template, input_path, TEMP_DIR, target_format, template_style
        )
        
        background_tasks.add_task(cleanup_file, input_path)
        background_tasks.add_task(cleanup_file, output_path)

        final_filename = f"converted_{Path(file.filename).stem}{target_format}"
        return FileResponse(path=str(output_path), filename=final_filename)
    
    except Exception as e:
        logger.error(f"Gagal konversi file {file.filename}: {e}")
        if input_path.exists(): input_path.unlink()
        raise HTTPException(status_code=400, detail=str(e))
app.mount("/", StaticFiles(directory="../frontend", html=True), name="frontend")